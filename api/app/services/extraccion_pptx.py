# ==========================================================
# Extracción de .pptx a texto (Generador de prompts — Unidad de programación)
# ==========================================================
#
# Un marcador "### Diapositiva N" por diapositiva, mismo criterio que ya se
# validó dos veces con datos reales en el prototipo de esta sesión
# (generar_prompt.py / atmosfera.txt): la instrucción del prompt le pide a la
# IA que revise diapositiva a diapositiva antes de responder, así que el
# marcador tiene que ser inequívoco y constante.
#
# No recorre shapes agrupados (GroupShape) -- mismo criterio de alcance que
# services/extraccion_docx.py con encabezados/pies de página: cubre el caso
# real (texto y tablas sueltas en la diapositiva), no cada rincón posible del
# formato.

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from services.llm_client import transcribir_imagen

# Por debajo de esto, caracteres de texto extraído en la diapositiva, se
# intenta el fallback de visión -- pero solo si además hay una imagen en la
# diapositiva que transcribir (una diapositiva vacía de verdad, sin imagen,
# no tiene nada que la visión pueda rescatar).
_UMBRAL_CARACTERES_DIAPOSITIVA = 40


def _celda_a_texto(cell):

    return cell.text.strip()


def _tabla_a_markdown(table):

    filas_md = []

    for i, row in enumerate(table.rows):
        celdas = [_celda_a_texto(cell).replace("|", "\\|").replace("\n", " ") for cell in row.cells]
        filas_md.append("| " + " | ".join(celdas) + " |")

        if i == 0:
            filas_md.append("| " + " | ".join(["---"] * len(celdas)) + " |")

    return "\n".join(filas_md)


def extraer_texto_pptx(contenido_bytes):
    """Devuelve (texto, num_diapositivas, diapositivas_con_vision,
    diapositivas_vision_fallida). Una diapositiva con poco texto y una
    imagen grande (típico de una diapositiva escaneada pegada como imagen
    de fondo) se intenta releer con el modelo de visión del ia-server --
    aquí no hace falta renderizar la diapositiva entera, basta con la propia
    imagen incrustada."""

    presentacion = Presentation(io.BytesIO(contenido_bytes))
    bloques = []
    diapositivas_con_vision = []
    diapositivas_vision_fallida = []

    for i, diapositiva in enumerate(presentacion.slides, start=1):

        partes = []
        imagen_mayor = None

        for shape in diapositiva.shapes:

            if shape.has_text_frame:
                texto = shape.text_frame.text.strip()
                if texto:
                    partes.append(texto)

            elif shape.has_table:
                partes.append(_tabla_a_markdown(shape.table))

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if imagen_mayor is None or (shape.width * shape.height) > (imagen_mayor.width * imagen_mayor.height):
                    imagen_mayor = shape

        if diapositiva.has_notes_slide:
            notas = diapositiva.notes_slide.notes_text_frame.text.strip()
            if notas:
                partes.append(f"(Notas del orador: {notas})")

        texto_diapositiva = "\n\n".join(partes)

        if len(texto_diapositiva) < _UMBRAL_CARACTERES_DIAPOSITIVA and imagen_mayor is not None:
            try:
                imagen = imagen_mayor.image
                resultado = transcribir_imagen(imagen.blob, mime_type=imagen.content_type)
            except Exception:
                resultado = None

            if resultado:
                texto_diapositiva = (texto_diapositiva + "\n\n" + resultado).strip()
                diapositivas_con_vision.append(i)
            else:
                diapositivas_vision_fallida.append(i)

        bloques.append(f"### Diapositiva {i}\n{texto_diapositiva}")

    return "\n\n".join(bloques), len(presentacion.slides), diapositivas_con_vision, diapositivas_vision_fallida
